from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
import uuid
import os
# Removed requests dependency - not needed for core functionality

# ---- CONFIG ----

# Use absolute path for the static directory
STATIC_PATH = os.path.abspath("generated_files")
os.makedirs(STATIC_PATH, exist_ok=True)

# ---- FASTAPI APP ----
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],   # Adjust for production!
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---- PYDANTIC MODEL ----
class SeminarKitRequest(BaseModel):
    topic: str
    name: str
    roll: str
    college: str
    semester: str
    branch: str

# ---- HELPER: GPT GENERATION ----
def gpt_generate(prompt, model="gpt-3.5-turbo", temperature=0.6, max_tokens=800):
    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo-0125",
            temperature=temperature,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=max_tokens,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"GPT Generation Error: {e}")
        return f"Error generating content for: {prompt[:50]}..."

# Removed image functionality to avoid external dependencies

# ---- HELPER: ENHANCED PPTX GENERATION ----
def generate_enhanced_pptx(topic, sections, student_info):
    try:
        file_id = str(uuid.uuid4())
        filename = f"seminar_{file_id}.pptx"
        file_path = os.path.join(STATIC_PATH, filename)
        
        prs = Presentation()
        
        # ---- SLIDE 1: TITLE SLIDE ----
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = topic.upper()
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)  # Blue color
        
        # Subtitle with student info
        subtitle_shape = slide.placeholders[1]
        subtitle_text = f"""Presented by: {student_info['name']}
Roll No: {student_info['roll']}
{student_info['college']}
Semester: {student_info['semester']} | Branch: {student_info['branch']}

Academic Seminar Presentation"""
        
        subtitle_shape.text = subtitle_text
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # ---- SLIDE 2: AGENDA ----
        bullet_slide_layout = prs.slide_layouts[1]  # Content with bullets layout
        agenda_slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Agenda title
        agenda_title = agenda_slide.shapes.title
        agenda_title.text = "PRESENTATION AGENDA"
        agenda_title.text_frame.paragraphs[0].font.size = Pt(36)
        agenda_title.text_frame.paragraphs[0].font.bold = True
        agenda_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)
        
        # Agenda content
        agenda_content = agenda_slide.placeholders[1]
        agenda_text = ""
        for i, section in enumerate(sections, 1):
            agenda_text += f"{i}. {section['title']}\n"
        
        agenda_content.text = agenda_text
        for paragraph in agenda_content.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.space_after = Pt(12)
        
        # ---- CONTENT SLIDES ----
        for section in sections:
            # Create slide for each section
            content_slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Section title
            section_title = content_slide.shapes.title
            section_title.text = section['title'].upper()
            section_title.text_frame.paragraphs[0].font.size = Pt(32)
            section_title.text_frame.paragraphs[0].font.bold = True
            section_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)
            
            # Section content
            section_content = content_slide.placeholders[1]
            content_text = ""
            
            # Add main points
            for point in section['points']:
                content_text += f"• {point}\n"
            
            # Add sub-points if available
            if 'sub_points' in section:
                content_text += "\nKey Details:\n"
                for sub_point in section['sub_points']:
                    content_text += f"  ◦ {sub_point}\n"
            
            section_content.text = content_text
            
            # Format content text
            for paragraph in section_content.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.space_after = Pt(8)
                if paragraph.text.startswith('•'):
                    paragraph.font.bold = True
                elif paragraph.text.startswith('  ◦'):
                    paragraph.font.size = Pt(18)
        
        # ---- CONCLUSION SLIDE ----
        conclusion_slide = prs.slides.add_slide(bullet_slide_layout)
        
        conclusion_title = conclusion_slide.shapes.title
        conclusion_title.text = "CONCLUSION & FUTURE SCOPE"
        conclusion_title.text_frame.paragraphs[0].font.size = Pt(32)
        conclusion_title.text_frame.paragraphs[0].font.bold = True
        conclusion_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)
        
        conclusion_content = conclusion_slide.placeholders[1]
        conclusion_text = f"""• {topic} represents a significant advancement in technology
• Widespread applications across multiple industries
• Continuous research and development ongoing
• Future implementation will revolutionize current practices
• Potential for further innovation and improvement

Thank you for your attention!
Questions & Discussion"""
        
        conclusion_content.text = conclusion_text
        for paragraph in conclusion_content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.space_after = Pt(8)
        
        # Save the presentation
        prs.save(file_path)
        
        # Verify file was created
        if not os.path.exists(file_path):
            raise Exception(f"PPTX file was not created at {file_path}")
        
        print(f"Enhanced PPTX created successfully at: {file_path}")
        return filename
        
    except Exception as e:
        print(f"Error creating enhanced PPTX: {e}")
        raise Exception(f"Failed to create enhanced PPTX: {str(e)}")

def generate_detailed_content(topic):
    """Generate detailed, structured content for the presentation"""
    
    # Generate comprehensive content structure
    structure_prompt = f"""Create a detailed seminar presentation structure for the topic "{topic}".
    
    Generate exactly 6 main sections with the following format:
    
    Section 1: Introduction
    - Brief overview
    - Importance and relevance
    - Scope of discussion
    
    Section 2: Historical Background / Literature Review
    - Origins and development
    - Key milestones
    - Previous research
    
    Section 3: Core Concepts / Theoretical Framework
    - Fundamental principles
    - Key definitions
    - Technical aspects
    
    Section 4: Applications / Implementation
    - Real-world applications
    - Use cases
    - Industry adoption
    
    Section 5: Challenges and Solutions
    - Current limitations
    - Technical challenges
    - Proposed solutions
    
    Section 6: Future Trends / Research Directions
    - Emerging trends
    - Future possibilities
    - Research opportunities
    
    For each section, provide:
    - 3-4 main bullet points
    - 2-3 detailed sub-points for each main point
    
    Make the content technical and suitable for academic presentation."""
    
    try:
        detailed_content = gpt_generate(structure_prompt, max_tokens=1500)
        
        # Parse the generated content into structured sections
        sections = []
        current_section = None
        
        lines = detailed_content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('Section') and ':' in line:
                if current_section:
                    sections.append(current_section)
                
                section_title = line.split(':', 1)[1].strip()
                current_section = {
                    'title': section_title,
                    'points': [],
                    'sub_points': []
                }
            elif line.startswith('-') and current_section:
                point = line.lstrip('- ').strip()
                if point:
                    current_section['points'].append(point)
            elif line.startswith('•') and current_section:
                point = line.lstrip('• ').strip()
                if point:
                    current_section['points'].append(point)
        
        # Add the last section
        if current_section:
            sections.append(current_section)
        
        # Fallback structured content if parsing fails
        if len(sections) < 4:
            sections = [
                {
                    'title': 'Introduction and Overview',
                    'points': [
                        f'Definition and scope of {topic}',
                        'Importance in current technological landscape',
                        'Objectives of this presentation',
                        'Expected outcomes and benefits'
                    ]
                },
                {
                    'title': 'Literature Review and Background',
                    'points': [
                        'Historical development and evolution',
                        'Key research contributions',
                        'Current state of technology',
                        'Gaps in existing knowledge'
                    ]
                },
                {
                    'title': 'Core Concepts and Methodology',
                    'points': [
                        'Fundamental principles and theories',
                        'Technical architecture and design',
                        'Implementation methodologies',
                        'Performance metrics and evaluation'
                    ]
                },
                {
                    'title': 'Applications and Case Studies',
                    'points': [
                        'Industry applications and use cases',
                        'Real-world implementation examples',
                        'Success stories and best practices',
                        'Comparative analysis with alternatives'
                    ]
                },
                {
                    'title': 'Challenges and Limitations',
                    'points': [
                        'Technical challenges and constraints',
                        'Implementation barriers',
                        'Cost and resource considerations',
                        'Scalability and performance issues'
                    ]
                },
                {
                    'title': 'Future Scope and Recommendations',
                    'points': [
                        'Emerging trends and opportunities',
                        'Research directions and possibilities',
                        'Recommendations for implementation',
                        'Expected future developments'
                    ]
                }
            ]
        
        return sections
        
    except Exception as e:
        print(f"Error generating detailed content: {e}")
        # Return fallback content
        return [
            {
                'title': 'Introduction',
                'points': [f'Overview of {topic}', 'Significance and applications', 'Presentation objectives']
            },
            {
                'title': 'Technical Details',
                'points': ['Core concepts', 'Implementation methods', 'Key features']
            },
            {
                'title': 'Applications',
                'points': ['Industry use cases', 'Practical examples', 'Benefits and advantages']
            },
            {
                'title': 'Future Scope',
                'points': ['Emerging trends', 'Research opportunities', 'Potential developments']
            }
        ]

def generate_docx(topic, sections, student_info, qna):
    try:
        file_id = str(uuid.uuid4())
        filename = f"seminar_{file_id}.docx"
        file_path = os.path.join(STATIC_PATH, filename)
        
        doc = Document()
        
        # Cover page
        doc.add_heading(topic.upper(), 0)
        doc.add_paragraph(f"Seminar Report")
        doc.add_paragraph("")
        doc.add_paragraph(f"Submitted by: {student_info['name']}")
        doc.add_paragraph(f"Roll Number: {student_info['roll']}")
        doc.add_paragraph(f"College: {student_info['college']}")
        doc.add_paragraph(f"Semester: {student_info['semester']}")
        doc.add_paragraph(f"Branch: {student_info['branch']}")
        doc.add_page_break()
        
        # Table of Contents
        doc.add_heading("TABLE OF CONTENTS", level=1)
        for i, section in enumerate(sections, 1):
            doc.add_paragraph(f"{i}. {section['title']}")
        doc.add_paragraph(f"{len(sections)+1}. Questions & Answers")
        doc.add_paragraph(f"{len(sections)+2}. Conclusion")
        doc.add_page_break()
        
        # Content sections
        for i, section in enumerate(sections, 1):
            doc.add_heading(f"{i}. {section['title']}", level=1)
            
            for point in section['points']:
                doc.add_paragraph(f"• {point}")
                
            if 'sub_points' in section and section['sub_points']:
                doc.add_paragraph("Detailed Information:")
                for sub_point in section['sub_points']:
                    doc.add_paragraph(f"  ◦ {sub_point}")
            
            doc.add_paragraph("")  # Add spacing
        
        # Q&A section
        doc.add_heading(f"{len(sections)+1}. Questions & Answers", level=1)
        for i, qa in enumerate(qna, 1):
            doc.add_paragraph(f"Q{i}: {qa['question']}", style="Heading 3")
            doc.add_paragraph(f"Answer: {qa['answer']}")
            doc.add_paragraph("")
        
        # Conclusion
        doc.add_heading(f"{len(sections)+2}. Conclusion", level=1)
        doc.add_paragraph(f"This seminar on {topic} has provided comprehensive insights into the subject matter. The discussion covered various aspects including technical details, applications, and future scope. The research and analysis presented demonstrate the significance of {topic} in current technological advancements.")
        
        # Save the document
        doc.save(file_path)
        
        # Verify file was created
        if not os.path.exists(file_path):
            raise Exception(f"DOCX file was not created at {file_path}")
        
        file_size = os.path.getsize(file_path)
        print(f"Enhanced DOCX created successfully at: {file_path} (Size: {file_size} bytes)")
        return filename
        
    except Exception as e:
        print(f"Error creating enhanced DOCX: {e}")
        raise Exception(f"Failed to create enhanced DOCX: {str(e)}")

# ---- API ENDPOINT ----
@app.post("/generate-seminar-kit")
async def generate_kit(body: SeminarKitRequest):
    try:
        student_info = {
            "name": body.name,
            "roll": body.roll,
            "college": body.college,
            "semester": body.semester,
            "branch": body.branch
        }
        
        print(f"Generating enhanced seminar kit for topic: {body.topic}")
        
        # Generate detailed content structure
        sections = generate_detailed_content(body.topic)
        print(f"Generated {len(sections)} detailed sections")
        
        # Generate Q&A
        qna_prompt = f"""Generate exactly 8 comprehensive questions and detailed answers for a student seminar presentation on '{body.topic}'. 
        
        Include questions about:
        - Basic concepts and definitions
        - Technical implementation
        - Applications and use cases  
        - Advantages and disadvantages
        - Future scope and trends
        - Comparison with alternatives
        - Challenges and solutions
        - Research opportunities
        
        Format each as:
        Q: [detailed question]
        A: [comprehensive answer in 2-3 sentences]
        
        Make answers technical and suitable for academic presentation."""

        qna_text = gpt_generate(qna_prompt, max_tokens=1200)
        
        # Parse Q&A with improved logic
        qna_pairs = []
        current_q = ""
        current_a = ""
        
        for line in qna_text.split('\n'):
            line = line.strip()
            if line.startswith('Q:') or line.startswith('Question'):
                if current_q and current_a:
                    qna_pairs.append({"question": current_q, "answer": current_a})
                current_q = line.replace('Q:', '').replace('Question:', '').strip()
                current_a = ""
            elif line.startswith('A:') or line.startswith('Answer'):
                current_a = line.replace('A:', '').replace('Answer:', '').strip()
            elif current_a and line and not line.startswith('Q:'):
                current_a += " " + line
        
        # Add the last Q&A pair
        if current_q and current_a:
            qna_pairs.append({"question": current_q, "answer": current_a})
        
        # Enhanced fallback Q&A if parsing failed
        if len(qna_pairs) < 5:
            qna_pairs = [
                {"question": f"What is {body.topic} and why is it important?", 
                 "answer": f"{body.topic} is a significant technological advancement that plays a crucial role in modern applications. Its importance lies in its ability to solve complex problems and improve efficiency in various domains."},
                {"question": f"What are the key technical components of {body.topic}?", 
                 "answer": f"The main technical components include core algorithms, implementation frameworks, and supporting infrastructure. These components work together to provide comprehensive functionality."},
                {"question": f"What are the real-world applications of {body.topic}?", 
                 "answer": f"{body.topic} finds applications in multiple industries including healthcare, finance, manufacturing, and telecommunications. It enables automation and optimization of complex processes."},
                {"question": f"What are the main advantages and limitations of {body.topic}?", 
                 "answer": f"Key advantages include improved efficiency, scalability, and cost-effectiveness. However, limitations may include implementation complexity, resource requirements, and technical constraints."},
                {"question": f"What does the future hold for {body.topic}?", 
                 "answer": f"Future developments in {body.topic} are expected to focus on enhanced performance, broader applications, and integration with emerging technologies. Research continues to address current limitations."},
                {"question": f"How does {body.topic} compare with alternative approaches?", 
                 "answer": f"Compared to traditional methods, {body.topic} offers superior performance in terms of speed, accuracy, and scalability. However, the choice depends on specific requirements and constraints."},
                {"question": f"What are the main challenges in implementing {body.topic}?", 
                 "answer": f"Implementation challenges include technical complexity, resource allocation, skill requirements, and integration with existing systems. Proper planning and expertise are essential for success."},
                {"question": f"What research opportunities exist in {body.topic}?", 
                 "answer": f"Research opportunities include algorithm optimization, new application domains, performance enhancement, and addressing current limitations. Interdisciplinary collaboration can lead to breakthrough innovations."}
            ]
        
        print(f"Generated {len(qna_pairs)} Q&A pairs")

        # Generate files with enhanced content
        pptx_filename = generate_enhanced_pptx(body.topic, sections, student_info)
        docx_filename = generate_docx(body.topic, sections, student_info, qna_pairs)
        
        # Generate enhanced QnA text file
        qna_filename = f"qna_{uuid.uuid4()}.txt"
        qna_file_path = os.path.join(STATIC_PATH, qna_filename)
        
        with open(qna_file_path, "w", encoding="utf-8") as f:
            f.write(f"COMPREHENSIVE QUESTIONS & ANSWERS\n")
            f.write(f"Topic: {body.topic}\n")
            f.write(f"Prepared for: {student_info['name']} ({student_info['roll']})\n")
            f.write(f"College: {student_info['college']}\n")
            f.write(f"Semester: {student_info['semester']} | Branch: {student_info['branch']}\n")
            f.write("="*80 + "\n\n")
            
            for i, qa in enumerate(qna_pairs, 1):
                f.write(f"QUESTION {i}:\n{qa['question']}\n\n")
                f.write(f"ANSWER {i}:\n{qa['answer']}\n")
                f.write("-"*60 + "\n\n")

        return {
            "success": True,
            "message": "Enhanced seminar kit generated successfully!",
            "pptx_url": f"http://localhost:8000/download/{pptx_filename}",
            "docx_url": f"http://localhost:8000/download/{docx_filename}",
            "qna_url": f"http://localhost:8000/download/{qna_filename}",
            "details": {
                "sections_generated": len(sections),
                "questions_generated": len(qna_pairs),
                "presentation_slides": len(sections) + 3  # Title + Agenda + Conclusion
            }
        }
        
    except Exception as e:
        print(f"Error in generate_kit: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate enhanced seminar kit: {str(e)}")

# ---- FILE SERVING ENDPOINT ----
from fastapi.responses import FileResponse

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(STATIC_PATH, filename)
    
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        raise HTTPException(status_code=404, detail="File not found")
    
    # Check file size
    file_size = os.path.getsize(file_path)
    if file_size == 0:
        print(f"File is empty: {file_path}")
        raise HTTPException(status_code=500, detail="Generated file is empty")
    
    print(f"Serving file: {file_path} (Size: {file_size} bytes)")
    
    # Determine media type based on file extension
    if filename.endswith('.docx'):
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif filename.endswith('.pptx'):
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif filename.endswith('.txt'):
        media_type = "text/plain"
    else:
        media_type = "application/octet-stream"
    
    return FileResponse(
        path=file_path, 
        filename=filename, 
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# Add a health check endpoint
@app.get("/health")
async def health_check():
    return {"status": "healthy", "static_path": STATIC_PATH}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
